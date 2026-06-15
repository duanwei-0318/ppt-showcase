"""
PPT Processor - Extract slides and auto-categorize products from PPTX files.
"""
import os
from PIL import Image, ImageDraw, ImageFont
import re
import json
import uuid
from datetime import datetime
from pptx import Presentation


def extract_text_from_slide(slide):
    """Extract all text from a slide, returning structured content."""
    title = ""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)
    if texts:
        title = texts[0]
    return title, texts


def extract_images_from_slide(slide, output_dir, slide_index):
    """Extract images from a slide and save them to disk."""
    image_paths = []
    img_count = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            try:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                filename = f"slide_{slide_index}_img_{img_count}.{ext}"
                filepath = os.path.join(output_dir, filename)
                with open(filepath, "wb") as f:
                    f.write(image.blob)
                image_paths.append(f"static/images/{filename}")
                img_count += 1
            except Exception:
                pass
    return image_paths


CATEGORY_KEYWORDS = {
    "电子产品": ["手机", "电脑", "平板", "耳机", "相机", "电子", "数码", "智能", "充电", "电池", "屏幕", "芯片", "设备"],
    "家居生活": ["家具", "家居", "床", "桌", "椅", "沙发", "灯", "装饰", "厨房", "浴室", "收纳", "家纺", "窗帘", "地毯"],
    "服装配饰": ["服装", "衣服", "衬衫", "裤子", "裙子", "鞋", "包", "配饰", "手表", "首饰", "帽子", "围巾", "皮带"],
    "食品饮料": ["食品", "饮料", "零食", "水果", "茶", "咖啡", "酒", "巧克力", "饼干", "面包", "蛋糕", "牛奶", "水"],
    "运动户外": ["运动", "健身", "跑步", "瑜伽", "户外", "登山", "露营", "骑行", "游泳", "球", "鞋", "装备", "器材"],
    "办公用品": ["办公", "文具", "笔", "纸", "文件夹", "打印", "电脑", "桌", "椅子", "书", "本", "胶带", "剪刀"],
    "美妆护肤": ["美妆", "护肤", "化妆", "面膜", "精华", "面霜", "口红", "眼影", "香水", "洗发", "沐浴", "防晒"],
    "图书教育": ["图书", "书", "教育", "学习", "课程", "教材", "培训", "知识", "阅读", "小说", "文学", "科学"]
}


def auto_categorize(texts, existing_categories):
    """Auto-categorize slide content based on keyword matching."""
    scores = {cat: 0 for cat in CATEGORY_KEYWORDS}
    all_text = " ".join(texts).lower()

    for cat, keywords in CATEGORY_KEYWORDS.items():
        for kw in keywords:
            if kw.lower() in all_text:
                scores[cat] += 1

    best = max(scores, key=scores.get)
    if scores[best] > 0:
        return best

    # Check against existing custom categories
    for cat in existing_categories:
        cat_name = cat["name"].lower()
        if cat_name in all_text:
            return cat["name"]

    return "未分类"


def generate_placeholder(output_dir, slide_index, title):
    cat_colors = {
        "电子产品": (9, 132, 227), "家居生活": (108, 92, 231), "服装配饰": (225, 112, 85),
        "食品饮料": (0, 184, 148), "运动户外": (0, 206, 201), "办公用品": (116, 185, 255),
        "美妆护肤": (253, 121, 168), "图书教育": (162, 155, 254),
    }
    r, g, b = cat_colors.get(title[:4], 45, 52, 54)
    img = Image.new("RGB", (400, 300), (r, g, b))
    draw = ImageDraw.Draw(img)
    # Draw centered diamond shape
    cx, cy = 200, 120
    s = 40
    draw.polygon([(cx, cy-s), (cx+s, cy), (cx, cy+s), (cx-s, cy)], fill=(255, 255, 255, 40))
    # Draw product icon
    draw.text((cx - 12, cy - 10), "◆", fill=(255, 255, 255, 200))
    # Draw title text
    lines = [title[i:i+10] for i in range(0, len(title), 10)]
    y = 210
    for line in lines[:2]:
        tw = len(line) * 8
        draw.text(((400-tw)/2, y), line, fill=(255, 255, 255, 220))
        y += 22
    filename = f"slide_{slide_index}_placeholder.png"
    filepath = os.path.join(output_dir, filename)
    img.save(filepath, "PNG")
    return f"static/images/{filename}"


def extract_price(texts):
    """Try to extract price information from slide text."""
    price_patterns = [
        r'[¥￥](\d+(?:[.,]\d+)?)',
        r'价格[：:]\s*(\d+(?:[.,]\d+)?)',
        r'售价[：:]\s*(\d+(?:[.,]\d+)?)',
        r'(\d+(?:[.,]\d+)?)\s*元',
    ]
    all_text = " ".join(texts)
    for pattern in price_patterns:
        match = re.search(pattern, all_text)
        if match:
            return float(match.group(1).replace(",", ""))
    return None


def process_pptx(pptx_path, output_dir, existing_categories):
    """Process a PPTX file and extract all slides as products."""
    prs = Presentation(pptx_path)
    products = []
    pptx_filename = os.path.basename(pptx_path)

    for i, slide in enumerate(prs.slides):
        title, texts = extract_text_from_slide(slide)
        image_paths = extract_images_from_slide(slide, output_dir, i)

        category = auto_categorize(texts, existing_categories)
        price = extract_price(texts)

        product = {
            "id": str(uuid.uuid4()),
            "title": title or f"产品 {i+1}",
            "description": "\n".join(texts[1:]) if len(texts) > 1 else "",
            "price": price if price else 0,
            "category": category,
            "images": image_paths,
            "source_file": pptx_filename,
            "slide_index": i,
            "created_at": datetime.now().isoformat(),
            "all_text": " ".join(texts)
        }
        products.append(product)

    return products


if __name__ == "__main__":
    print("PPT Processor module loaded successfully")
